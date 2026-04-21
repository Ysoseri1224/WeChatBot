import os
import re
import time
import json
import base64
import logging
import shutil
import sqlite3
import subprocess
import xml.etree.ElementTree as ET
from datetime import datetime, timezone, timedelta
from pathlib import Path
import queue
from queue import Empty
from threading import Thread
from dotenv import load_dotenv
from openai import OpenAI
from wcferry import Wcf, WxMsg

load_dotenv()

logging.basicConfig(
    level=logging.DEBUG,
    format="%(asctime)s [%(levelname)s] %(message)s",
    handlers=[
        logging.StreamHandler(),
        logging.FileHandler("bot.log", encoding="utf-8"),
    ],
)
LOG = logging.getLogger("WeChatBot")

GROUP_TRIGGER_PREFIX = os.getenv("GROUP_TRIGGER_PREFIX", "").strip()
PRIVATE_CHAT_REPLY = os.getenv("PRIVATE_CHAT_REPLY", "true").lower() == "true"
GROUP_WHITELIST = [g.strip() for g in os.getenv("GROUP_WHITELIST", "").split(",") if g.strip()]
SYSTEM_PROMPT = os.getenv("SYSTEM_PROMPT", "你是一个智能助手，请用中文简洁地回答用户的问题。")
WECHAT_PATH = os.getenv("WECHAT_PATH", "").strip() or None

AI_PROVIDER = os.getenv("AI_PROVIDER", "deepseek").lower().strip()

_PROVIDER_CONFIGS = {
    "gpt":      (os.getenv("GPT_API_KEY"),      os.getenv("GPT_BASE_URL",      "https://api.openai.com/v1"),   os.getenv("GPT_MODEL",      "gpt-4o-mini")),
    "deepseek": (os.getenv("DEEPSEEK_API_KEY"), os.getenv("DEEPSEEK_BASE_URL", "https://api.deepseek.com"),    os.getenv("DEEPSEEK_MODEL", "deepseek-chat")),
    "claude":   (os.getenv("CLAUDE_API_KEY"),   os.getenv("CLAUDE_BASE_URL",   "https://api.anthropic.com/v1"),os.getenv("CLAUDE_MODEL",   "claude-3-5-haiku-20241022")),
    "minimax":  (os.getenv("MINIMAX_API_KEY"),  os.getenv("MINIMAX_BASE_URL",  "https://api.minimax.io/v1"),   os.getenv("MINIMAX_MODEL",  "MiniMax-M2.5-highspeed")),
}

if AI_PROVIDER not in _PROVIDER_CONFIGS:
    raise ValueError(f"未知 AI_PROVIDER: {AI_PROVIDER}，可选: {list(_PROVIDER_CONFIGS.keys())}")

_api_key, _base_url, AI_MODEL = _PROVIDER_CONFIGS[AI_PROVIDER]

client = OpenAI(api_key=_api_key, base_url=_base_url)

conversation_history: dict = {}
group_whitelist_ids: set = set()
group_context_buffer: dict = {}   # roomid -> list of "昵称: 内容" 字符串
group_pending_media: dict = {}    # roomid -> {"type": "image", "msg": WxMsg}
_global_wcf = None                # 全局wcf引用，供文件监控线程使用
_file_watcher_started = False     # 文件监控线程是否已启动
GROUP_CONTEXT_LIMIT = int(os.getenv("GROUP_CONTEXT_LIMIT", "50"))  # 最多保留条数

VISION_PROVIDERS = {"gpt", "claude"}  # 支持图片的供应商
WECHAT_FILE_DIR = os.getenv("WECHAT_FILE_DIR", r"D:\WeChat Files\wxid_amkb0miro4hf22\FileStorage\File")
MEMORY_DIR = Path(os.getenv("MEMORY_DIR", r"D:\Weixin\bot\memory"))
MEMORY_DIR.mkdir(parents=True, exist_ok=True)
FILE_SAVE_DIR = MEMORY_DIR / "files"
FILE_SAVE_DIR.mkdir(parents=True, exist_ok=True)
NOTES_DIR = MEMORY_DIR / "notes"
NOTES_DIR.mkdir(parents=True, exist_ok=True)
SCHEDULES_DB = MEMORY_DIR / "schedules.db"
SOFFICE_PATH = os.getenv("SOFFICE_PATH", r"C:\Program Files\LibreOffice\program\soffice.exe")

BEIJING_TZ = timezone(timedelta(hours=8))
pending_schedule_confirm: dict = {}  # roomid -> {content, times, ts}
_schedule_reminder_started = False
WEEKDAYS_ZH = "一二三四五六日"


_MEMORY_DISABLED_FILE = None  # 延迟初始化，依赖 MEMORY_DIR


def _memory_disabled_path() -> Path:
    return MEMORY_DIR / "disabled.json"


def memory_disabled_load() -> set:
    p = _memory_disabled_path()
    if p.exists():
        try:
            return set(json.loads(p.read_text(encoding="utf-8")))
        except Exception:
            pass
    return set()


def memory_disabled_save(disabled: set):
    _memory_disabled_path().write_text(json.dumps(sorted(disabled), ensure_ascii=False), encoding="utf-8")


def memory_disabled_toggle(filename: str) -> bool:
    """切换禁用状态，返回切换后是否为禁用"""
    disabled = memory_disabled_load()
    if filename in disabled:
        disabled.discard(filename)
        memory_disabled_save(disabled)
        return False
    else:
        disabled.add(filename)
        memory_disabled_save(disabled)
        return True


def memory_load_all() -> str:
    disabled = memory_disabled_load()
    files = sorted(MEMORY_DIR.glob("*.txt"))
    if not files:
        return ""
    parts = []
    for f in files:
        if f.name in disabled:
            continue
        parts.append(f"=== 记忆：{f.stem} ===\n{f.read_text(encoding='utf-8', errors='ignore').strip()}")
    return "\n\n".join(parts)


def memory_list() -> str:
    files = sorted(MEMORY_DIR.glob("*.txt"))
    if not files:
        return "当前没有任何记忆。"
    lines = []
    for i, f in enumerate(files, 1):
        size = f.stat().st_size
        mtime = datetime.fromtimestamp(f.stat().st_mtime).strftime("%m-%d %H:%M")
        lines.append(f"{i}. {f.name}  ({size}B, {mtime})")
    return "\n".join(lines)


def memory_save(name: str, content: str) -> str:
    safe = re.sub(r'[\\/:*?"<>|]', '_', name.strip())
    if not safe.endswith(".txt"):
        safe += ".txt"
    path = MEMORY_DIR / safe
    path.write_text(content.strip(), encoding="utf-8")
    return safe


def memory_delete(name: str) -> bool:
    safe = re.sub(r'[\\/:*?"<>|]', '_', name.strip())
    if not safe.endswith(".txt"):
        safe += ".txt"
    path = MEMORY_DIR / safe
    if path.exists():
        path.unlink()
        return True
    return False


def memory_read(name: str) -> str:
    safe = re.sub(r'[\\/:*?"<>|]', '_', name.strip())
    if not safe.endswith(".txt"):
        safe += ".txt"
    path = MEMORY_DIR / safe
    if path.exists():
        return path.read_text(encoding="utf-8", errors="ignore").strip()
    return ""


# ── 工具 ──────────────────────────────────────────────────────────────
def now_beijing():
    return datetime.now(BEIJING_TZ)


# ── 笔记 ──────────────────────────────────────────────────────────────
def note_save(name: str, content: str) -> str:
    safe = re.sub(r'[\\/:*?"<>|]', '_', name.strip())
    p = NOTES_DIR / f"{safe}.txt"
    p.write_text(content.strip(), encoding="utf-8")
    return safe


def note_list() -> str:
    files = sorted(NOTES_DIR.glob("*.txt"))
    if not files:
        return "暂无笔记。"
    lines = ["📝 笔记列表："]
    for i, f in enumerate(files, 1):
        mtime = datetime.fromtimestamp(f.stat().st_mtime).strftime("%m-%d %H:%M")
        lines.append(f"  {i}. {f.stem}  ({mtime})")
    return "\n".join(lines)


def note_read(name: str):
    safe = re.sub(r'[\\/:*?"<>|]', '_', name.strip())
    exact = NOTES_DIR / f"{safe}.txt"
    if exact.exists():
        return exact.read_text(encoding="utf-8", errors="replace"), safe
    candidates = list(NOTES_DIR.glob(f"*{safe}*.txt"))
    if not candidates:
        return None, None
    f = candidates[0]
    return f.read_text(encoding="utf-8", errors="replace"), f.stem


# ── 日程 SQLite ────────────────────────────────────────────────────────
def _init_schedules_db():
    conn = sqlite3.connect(str(SCHEDULES_DB))
    conn.execute("""
        CREATE TABLE IF NOT EXISTS schedules (
            id       INTEGER PRIMARY KEY AUTOINCREMENT,
            name     TEXT NOT NULL,
            content  TEXT,
            year     INTEGER,
            month    INTEGER,
            day      INTEGER,
            hour     INTEGER,
            minute   INTEGER,
            weekday  INTEGER,
            reminded_morning INTEGER DEFAULT 0,
            reminded_before  INTEGER DEFAULT 0,
            roomid   TEXT,
            created_at TEXT
        )
    """)
    conn.commit()
    conn.close()


def schedule_add(name, content, year, month, day, hour, minute, weekday, roomid):
    conn = sqlite3.connect(str(SCHEDULES_DB))
    conn.execute(
        "INSERT INTO schedules (name,content,year,month,day,hour,minute,weekday,roomid,created_at) "
        "VALUES (?,?,?,?,?,?,?,?,?,?)",
        (name, content, year, month, day, hour, minute, weekday, roomid, now_beijing().isoformat())
    )
    conn.commit()
    conn.close()


def schedule_list() -> str:
    conn = sqlite3.connect(str(SCHEDULES_DB))
    rows = conn.execute(
        "SELECT id,name,year,month,day,hour,minute,weekday FROM schedules "
        "ORDER BY year,month,day,hour,minute"
    ).fetchall()
    conn.close()
    if not rows:
        return "暂无日程。"
    lines = ["📅 日程列表："]
    for row in rows:
        id_, name, y, mo, d, h, mi, wd = row
        wd_str = f"周{WEEKDAYS_ZH[wd]}" if wd is not None else ""
        lines.append(f"  [{id_}] {name}  {y}/{mo:02d}/{d:02d} {h:02d}:{mi:02d} {wd_str}")
    return "\n".join(lines)


def schedule_query(q: str) -> str:
    conn = sqlite3.connect(str(SCHEDULES_DB))
    conds, params = [], []
    wd_map = {"一": 0, "二": 1, "三": 2, "四": 3, "五": 4, "六": 5, "日": 6, "天": 6}
    for ch, wd in wd_map.items():
        if f"周{ch}" in q or f"星期{ch}" in q:
            conds.append("weekday=?"); params.append(wd); break
    nums = re.findall(r'\d+', q)
    if len(nums) >= 3:
        y = int(nums[0]); y = y + 2000 if y < 100 else y
        conds += ["year=?", "month=?", "day=?"]; params += [y, int(nums[1]), int(nums[2])]
    elif len(nums) == 2:
        conds += ["month=?", "day=?"]; params += [int(nums[0]), int(nums[1])]
    elif len(nums) == 1:
        n = int(nums[0])
        if n > 2000:
            conds.append("year=?"); params.append(n)
        elif n > 12:
            conds.append("day=?"); params.append(n)
        else:
            conds.append("month=?"); params.append(n)
    if conds:
        sql = (f"SELECT id,name,year,month,day,hour,minute,weekday,content FROM schedules "
               f"WHERE ({' AND '.join(conds)}) OR (name LIKE ?) ORDER BY year,month,day,hour,minute")
        params.append(f"%{q}%")
    else:
        sql = ("SELECT id,name,year,month,day,hour,minute,weekday,content FROM schedules "
               "WHERE name LIKE ? ORDER BY year,month,day,hour,minute")
        params = [f"%{q}%"]
    rows = conn.execute(sql, params).fetchall()
    conn.close()
    if not rows:
        return "未找到匹配的日程。"
    lines = [f"🔍 找到 {len(rows)} 条日程："]
    for row in rows:
        id_, name, y, mo, d, h, mi, wd, content = row
        wd_str = f"周{WEEKDAYS_ZH[wd]}" if wd is not None else ""
        lines.append(f"  [{id_}] {name}  {y}/{mo:02d}/{d:02d} {h:02d}:{mi:02d} {wd_str}")
        if content:
            lines.append(f"       {content[:60]}")
    return "\n".join(lines)


# ── 时间自动检测 ───────────────────────────────────────────────────────
_TIME_KEYWORD_RE = re.compile(
    r'(今天|明天|后天|大后天|昨天|本周|这周|下周|上周|'
    r'周[一二三四五六日天]|星期[一二三四五六日天]|'
    r'\d{1,2}月\d{1,2}[号日]|\d{1,2}:\d{2}|'
    r'上午|下午|晚上|早上|下个月|这个月|本月|\d+点(钟|半)?|\d{4}年)'
)


def _has_time_keywords(text: str) -> bool:
    return bool(_TIME_KEYWORD_RE.search(text))


def _extract_times_with_ai(text: str) -> list:
    now_str = now_beijing().strftime(f"%Y年%m月%d日 %H:%M 周{WEEKDAYS_ZH[now_beijing().weekday()]}")
    prompt = (
        f"当前北京时间：{now_str}。\n"
        '从以下文本中提取所有具体时间点（含模糊表达如"下周一"、"明天下午"等），'
        '规范化为 YYYY-MM-DD HH:MM（时间不明确用 00:00），给出简短描述。\n'
        '只返回JSON数组，每项含 datetime 和 desc 字段，无时间点返回 []。\n\n'
        f"文本：\n{text}"
    )
    try:
        resp = client.chat.completions.create(
            model=AI_MODEL,
            messages=[{"role": "user", "content": prompt}],
            max_tokens=512,
            temperature=0,
        )
        raw = resp.choices[0].message.content.strip()
        m = re.search(r'\[.*\]', raw, re.DOTALL)
        if m:
            return json.loads(m.group())
    except Exception as e:
        LOG.warning(f"[时间提取] AI调用失败: {e}")
    return []


# ── 日程提醒线程 ───────────────────────────────────────────────────────
def _schedule_reminder_loop():
    global _global_wcf
    LOG.info("[日程提醒] 后台线程已启动")
    while True:
        try:
            wcf_ref = _global_wcf
            if wcf_ref and group_whitelist_ids:
                roomid = list(group_whitelist_ids)[0]
                now = now_beijing()
                conn = sqlite3.connect(str(SCHEDULES_DB))
                # 每天 07:00 提醒当日所有日程
                if now.hour == 7 and now.minute == 0:
                    rows = conn.execute(
                        "SELECT id,name,content,hour,minute,weekday FROM schedules "
                        "WHERE year=? AND month=? AND day=? AND reminded_morning=0",
                        (now.year, now.month, now.day)
                    ).fetchall()
                    for id_, name, content, h, mi, wd in rows:
                        wd_str = f"周{WEEKDAYS_ZH[wd]}" if wd is not None else ""
                        msg = f"📅 今日日程：{name}\n{h:02d}:{mi:02d} {wd_str}"
                        if content:
                            msg += f"\n\n{content}"
                        wcf_ref.send_text(msg, roomid)
                        conn.execute("UPDATE schedules SET reminded_morning=1 WHERE id=?", (id_,))
                    conn.commit()
                # 提前 1 小时提醒
                ahead = now + timedelta(hours=1)
                rows = conn.execute(
                    "SELECT id,name,content,weekday FROM schedules "
                    "WHERE year=? AND month=? AND day=? AND hour=? AND minute=? AND reminded_before=0",
                    (ahead.year, ahead.month, ahead.day, ahead.hour, ahead.minute)
                ).fetchall()
                for id_, name, content, wd in rows:
                    wd_str = f"周{WEEKDAYS_ZH[wd]}" if wd is not None else ""
                    msg = f"⏰ 1小时后：{name} {wd_str}\n{ahead.strftime('%H:%M')}"
                    if content:
                        msg += f"\n\n{content}"
                    wcf_ref.send_text(msg, roomid)
                    conn.execute("UPDATE schedules SET reminded_before=1 WHERE id=?", (id_,))
                conn.commit()
                conn.close()
        except Exception as e:
            LOG.warning(f"[日程提醒] 出错: {e}")
        time.sleep(60)


CONVERTIBLE_EXTS = {".docx", ".doc", ".xlsx", ".xls", ".pptx", ".ppt", ".pdf", ".txt", ".md"}
LIBRE_NEEDED_EXTS = {".doc", ".ppt", ".xls"}


def _soffice_convert(src: Path, out_dir: Path, fmt: str = "docx"):
    if not Path(SOFFICE_PATH).exists():
        LOG.warning(f"LibreOffice 未找到: {SOFFICE_PATH}")
        return None
    try:
        subprocess.run(
            [SOFFICE_PATH, "--headless", "--convert-to", fmt, "--outdir", str(out_dir), str(src)],
            timeout=60, capture_output=True
        )
        stem = src.stem
        result = out_dir / f"{stem}.{fmt}"
        return result if result.exists() else None
    except Exception as e:
        LOG.error(f"LibreOffice 转换失败: {e}")
        return None


def convert_to_markdown(filepath: str) -> tuple[str, bool]:
    """返回 (markdown文本, 是否经过转换)。转换=原格式非原生支持"""
    p = Path(filepath)
    ext = p.suffix.lower()
    converted = False
    try:
        # .doc/.ppt/.xls → LibreOffice 先转
        if ext in LIBRE_NEEDED_EXTS:
            converted = True
            fmt_map = {".doc": "docx", ".ppt": "pptx", ".xls": "xlsx"}
            target_fmt = fmt_map[ext]
            new_p = _soffice_convert(p, p.parent, target_fmt)
            if new_p is None:
                return "", converted
            p = new_p
            ext = p.suffix.lower()

        if ext in (".txt", ".md"):
            return p.read_text(encoding="utf-8", errors="ignore"), converted

        elif ext == ".docx":
            import docx
            doc = docx.Document(str(p))
            lines = []
            for para in doc.paragraphs:
                if not para.text.strip():
                    continue
                style = para.style.name.lower() if para.style else ""
                if "heading 1" in style:
                    lines.append(f"# {para.text}")
                elif "heading 2" in style:
                    lines.append(f"## {para.text}")
                elif "heading 3" in style:
                    lines.append(f"### {para.text}")
                else:
                    lines.append(para.text)
            for table in doc.tables:
                rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
                if rows:
                    lines.append("| " + " | ".join(rows[0]) + " |")
                    lines.append("|" + "---|" * len(rows[0]))
                    for row in rows[1:]:
                        lines.append("| " + " | ".join(row) + " |")
            return "\n".join(lines), converted

        elif ext in (".xlsx",):
            import openpyxl
            wb = openpyxl.load_workbook(str(p), read_only=True, data_only=True)
            lines = []
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                lines.append(f"## Sheet: {sheet}")
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) if c is not None else "" for c in row]
                    if any(cells):
                        lines.append("| " + " | ".join(cells) + " |")
            return "\n".join(lines), converted

        elif ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(str(p))
            lines = []
            for i, slide in enumerate(prs.slides, 1):
                lines.append(f"## 第 {i} 页")
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                lines.append(text)
            return "\n".join(lines), converted

        elif ext == ".pdf":
            import fitz
            doc = fitz.open(str(p))
            return "\n".join(page.get_text() for page in doc), converted

        else:
            return "", converted

    except Exception as e:
        LOG.error(f"文件转换失败 {filepath}: {e}")
        return "", converted


def extract_file_text(filepath: str) -> str:
    text, _ = convert_to_markdown(filepath)
    return text


_watched_files: set = set()  # 已处理过的文件路径集合
_file_convert_status: dict = {}  # stem -> {"status": "converting"/"done"/"error", "path": str, "md_path": str}
_new_files_pending: list = []   # 断线期间检测到的新文件，重连后通知用户

# ── 重连后从微信数据库恢复文件 ──
_last_recover_ts_file = MEMORY_DIR / ".last_recover_ts"


def _extract_filename_from_xml(xml_str):
    """从 type=49 消息的 XML 中提取文件名（尝试多个标签）"""
    if not xml_str:
        return ""
    s = xml_str if isinstance(xml_str, str) else str(xml_str)
    # 优先从 <filename> 提取（最准确）
    m = re.search(r'<filename>([^<]+)</filename>', s)
    if m:
        return m.group(1).strip()
    # 其次 <title>
    m = re.search(r'<title>([^<]+)</title>', s)
    if m:
        name = m.group(1).strip()
        if '.' in name:  # 确保是文件名而非链接标题
            return name
    # 最后 <sourcedisplayname>
    m = re.search(r'<sourcedisplayname>([^<]+)</sourcedisplayname>', s)
    if m:
        return m.group(1).strip()
    return ""


def _try_download_from_db(wcf, target_filename):
    """按需下载：查 DB 找匹配文件名的 type=49 消息，调 download_attach 下载。
    返回本地文件路径，失败返回 None。"""
    target_stem = Path(target_filename).stem.lower()
    target_ext = Path(target_filename).suffix.lower()

    try:
        dbs = wcf.get_dbs()
        msg_dbs = [d for d in dbs if d.startswith("MSG")]
        if not msg_dbs:
            LOG.warning("[按需下载] 未找到 MSG 数据库")
            return None

        # 查最近1小时的 type=49 消息
        min_ts = int(time.time()) - 3600
        candidates = []

        for db in msg_dbs:
            try:
                tables = wcf.get_tables(db)
                table_names = [t.get("name", "") for t in tables] if isinstance(tables, list) else []
                for tbl in table_names:
                    if not tbl.startswith("MSG"):
                        continue
                    sql = (
                        f"SELECT localId, MsgSvrID, StrTalker, StrContent, CreateTime "
                        f"FROM {tbl} WHERE Type=49 AND CreateTime>{min_ts} "
                        f"ORDER BY CreateTime DESC LIMIT 50"
                    )
                    try:
                        rows = wcf.query_sql(db, sql)
                        for row in rows:
                            raw = row.get("StrContent", "")
                            if isinstance(raw, bytes):
                                try:
                                    raw = raw.decode("utf-8")
                                except Exception:
                                    raw = raw.decode("utf-8", errors="replace")
                            fname = _extract_filename_from_xml(raw)
                            if not fname:
                                continue
                            f_stem = Path(fname).stem.lower()
                            f_ext = Path(fname).suffix.lower()
                            # 匹配：完全匹配、前缀匹配、包含匹配
                            if (f_stem == target_stem or
                                f_stem.startswith(target_stem) or
                                target_stem in f_stem):
                                if not target_ext or f_ext == target_ext:
                                    candidates.append({
                                        "msg_id": row.get("MsgSvrID", 0) or row.get("localId", 0),
                                        "filename": fname,
                                        "create_time": row.get("CreateTime", 0),
                                    })
                    except Exception as e:
                        LOG.debug(f"[按需下载] 查询 {db}.{tbl} 失败: {e}")
            except Exception as e:
                LOG.debug(f"[按需下载] 读取 {db} 表结构失败: {e}")

        if not candidates:
            LOG.info(f"[按需下载] DB中未找到匹配 {target_filename!r} 的文件消息")
            return None

        # 取最新的一条
        candidates.sort(key=lambda x: x["create_time"], reverse=True)
        best = candidates[0]
        msg_id = best["msg_id"]
        filename = best["filename"]
        LOG.info(f"[按需下载] 找到匹配: {filename} (msg_id={msg_id})，正在下载...")

        ret = wcf.download_attach(msg_id, "", "")
        LOG.info(f"[按需下载] download_attach 返回: {ret}")

        # 等待文件出现（最多30秒）
        stem = Path(filename).stem
        for _ in range(15):
            time.sleep(2)
            for root_dir, dirs, files in os.walk(WECHAT_FILE_DIR):
                for f in files:
                    if f == filename or (Path(f).stem.startswith(stem) and Path(f).suffix.lower() == Path(filename).suffix.lower()):
                        candidate_path = os.path.join(root_dir, f)
                        try:
                            if time.time() - os.path.getmtime(candidate_path) < 300:
                                LOG.info(f"[按需下载] 文件已落盘: {candidate_path}")
                                return candidate_path
                        except OSError:
                            continue

        LOG.warning(f"[按需下载] 30秒内未找到文件: {filename}")
        return None

    except Exception as e:
        LOG.error(f"[按需下载] 失败: {e}", exc_info=True)
        return None


def _recover_files_from_db(wcf, group_ids):
    """重连后查询微信数据库，找到最近的文件消息并下载+转化"""
    try:
        # 读取上次恢复时间戳，避免重复处理
        last_ts = 0
        if _last_recover_ts_file.exists():
            try:
                last_ts = int(_last_recover_ts_file.read_text().strip())
            except Exception:
                pass
        # 至少查最近10分钟
        min_ts = max(last_ts, int(time.time()) - 600)

        dbs = wcf.get_dbs()
        LOG.info(f"[文件恢复] 可用数据库: {dbs}")

        # 找 MSG 相关的数据库
        msg_dbs = [d for d in dbs if d.startswith("MSG")]
        if not msg_dbs:
            LOG.warning("[文件恢复] 未找到 MSG 数据库")
            return

        file_msgs = []
        for db in msg_dbs:
            try:
                # 先查表结构
                tables = wcf.get_tables(db)
                table_names = [t.get("name", "") for t in tables] if isinstance(tables, list) else []
                LOG.info(f"[文件恢复] {db} 的表: {table_names}")

                # 尝试查 MSG 表中 type=49 的最近消息
                for tbl in table_names:
                    if not tbl.startswith("MSG"):
                        continue
                    sql = (
                        f"SELECT localId, MsgSvrID, Type, StrTalker, StrContent, "
                        f"CreateTime, BytesExtra "
                        f"FROM {tbl} WHERE Type=49 AND CreateTime>{min_ts} "
                        f"ORDER BY CreateTime DESC LIMIT 20"
                    )
                    try:
                        rows = wcf.query_sql(db, sql)
                        for row in rows:
                            talker = row.get("StrTalker", "")
                            if group_ids and talker not in group_ids:
                                continue
                            file_msgs.append({
                                "local_id": row.get("localId", 0),
                                "msg_svr_id": row.get("MsgSvrID", 0),
                                "roomid": talker,
                                "content": row.get("StrContent", ""),
                                "create_time": row.get("CreateTime", 0),
                                "bytes_extra": row.get("BytesExtra", b""),
                            })
                    except Exception as e:
                        LOG.debug(f"[文件恢复] 查询 {db}.{tbl} 失败: {e}")
            except Exception as e:
                LOG.debug(f"[文件恢复] 读取 {db} 表结构失败: {e}")

        if not file_msgs:
            LOG.info("[文件恢复] 数据库中未找到最近的文件消息")
            return

        LOG.info(f"[文件恢复] 找到 {len(file_msgs)} 条文件消息")

        # 更新时间戳
        _last_recover_ts_file.write_text(str(int(time.time())))

        for msg_info in file_msgs:
            raw_content = msg_info["content"]
            # DB中StrContent可能是bytes，需要解码
            if isinstance(raw_content, bytes):
                try:
                    raw_content = raw_content.decode("utf-8")
                except Exception:
                    raw_content = raw_content.decode("utf-8", errors="replace")
            LOG.info(f"[文件恢复] StrContent前200字: {str(raw_content)[:200]}")
            filename = _extract_filename_from_xml(raw_content)
            ext = Path(filename).suffix.lower() if filename else ""
            if ext not in CONVERTIBLE_EXTS:
                LOG.debug(f"[文件恢复] 跳过非可转化文件: {filename!r}")
                continue

            roomid = msg_info["roomid"]
            msg_id = msg_info["msg_svr_id"] or msg_info["local_id"]
            LOG.info(f"[文件恢复] 处理文件: {filename} (id={msg_id}, room={roomid})")

            try:
                wcf.send_text(f"⏳ 正在下载文件《{filename}》...", roomid)
                # download_attach: thumb 和 extra 传空串尝试
                ret = wcf.download_attach(msg_id, "", "")
                LOG.info(f"[文件恢复] download_attach 返回: {ret}")

                # 等待文件出现（最多60秒）
                found_path = None
                stem = Path(filename).stem
                for _ in range(30):
                    time.sleep(2)
                    for root_dir, dirs, files in os.walk(WECHAT_FILE_DIR):
                        for f in files:
                            if f == filename or (Path(f).stem.startswith(stem) and Path(f).suffix.lower() == Path(filename).suffix.lower()):
                                candidate = os.path.join(root_dir, f)
                                try:
                                    if time.time() - os.path.getmtime(candidate) < 300:
                                        found_path = candidate
                                        break
                                except OSError:
                                    continue
                        if found_path:
                            break
                    if found_path:
                        break

                if not found_path:
                    LOG.warning(f"[文件恢复] 超时未找到: {filename}")
                    wcf.send_text(f"❌ 文件《{filename}》下载超时", roomid)
                    continue

                LOG.info(f"[文件恢复] 找到文件: {found_path}")
                tmp_path = FILE_SAVE_DIR / Path(found_path).name
                shutil.copy2(found_path, tmp_path)
                md_text, was_converted = convert_to_markdown(str(tmp_path))
                if tmp_path.suffix.lower() not in (".md", ".txt"):
                    tmp_path.unlink(missing_ok=True)
                if not md_text:
                    wcf.send_text(f"❌ 文件《{filename}》解析失败", roomid)
                    continue
                save_path = FILE_SAVE_DIR / f"{Path(found_path).stem}.md"
                save_path.write_text(md_text, encoding="utf-8")
                _file_convert_status[Path(found_path).stem] = {
                    "status": "done", "path": found_path,
                    "md_path": str(save_path), "name": filename
                }
                LOG.info(f"[文件恢复] 完成: {save_path}，{len(md_text)} 字符")
                wcf.send_text(
                    f"✅ 文件《{filename}》已下载并转化（{len(md_text)} 字符）\n"
                    f"可用：读取文件 {Path(found_path).stem}",
                    roomid
                )
            except Exception as e:
                LOG.error(f"[文件恢复] 处理 {filename} 失败: {e}", exc_info=True)
                try:
                    wcf.send_text(f"❌ 文件《{filename}》恢复失败: {e}", roomid)
                except Exception:
                    pass

    except Exception as e:
        LOG.error(f"[文件恢复] 整体失败: {e}", exc_info=True)


def _convert_and_save(src: Path) -> Path:
    """把 src 文件转换为 Markdown 并保存到 FILE_SAVE_DIR，返回保存路径。"""
    stem = src.stem
    _file_convert_status[stem] = {"status": "converting", "path": str(src), "md_path": "", "name": src.name}
    tmp_path = FILE_SAVE_DIR / src.name
    shutil.copy2(str(src), tmp_path)
    md_text, _ = convert_to_markdown(str(tmp_path))
    if tmp_path.suffix.lower() not in (".md", ".txt"):
        try:
            tmp_path.unlink()
        except Exception:
            pass
    if not md_text:
        _file_convert_status[stem]["status"] = "error"
        raise ValueError(f"文件解析为空: {src.name}")
    save_path = FILE_SAVE_DIR / f"{stem}.md"
    save_path.write_text(md_text, encoding="utf-8")
    _file_convert_status[stem]["status"] = "done"
    _file_convert_status[stem]["md_path"] = str(save_path)
    LOG.info(f"[转换] 完成: {save_path}，{len(md_text)} 字符")
    return save_path


def _try_write_env(key: str, value: str):
    """尝试把 key=value 写回 .env 文件（存在则更新，不存在则追加）。"""
    env_path = Path(__file__).parent / ".env"
    if not env_path.exists():
        return
    try:
        lines = env_path.read_text(encoding="utf-8").splitlines(keepends=True)
        found = False
        new_lines = []
        for line in lines:
            if re.match(rf"^{re.escape(key)}\s*=", line):
                new_lines.append(f"{key}={value}\n")
                found = True
            else:
                new_lines.append(line)
        if not found:
            new_lines.append(f"{key}={value}\n")
        env_path.write_text("".join(new_lines), encoding="utf-8")
    except Exception as e:
        LOG.warning(f"[WebUI] 写 .env 失败: {e}")


def _file_watcher():
    """后台线程：轮询微信缓存目录，检测新文件并静默转换保存。
    转换状态记录到 _file_convert_status，新文件记录到 _new_files_pending。"""
    global _global_wcf
    LOG.info(f"[文件监控] 线程启动，监控目录: {WECHAT_FILE_DIR}")
    start_time = time.time()
    for root_dir, dirs, files in os.walk(WECHAT_FILE_DIR):
        for f in files:
            _watched_files.add(os.path.join(root_dir, f))
    LOG.info(f"[文件监控] 已索引 {len(_watched_files)} 个现有文件")

    while True:
        try:
            time.sleep(5)
            for root_dir, dirs, files in os.walk(WECHAT_FILE_DIR):
                for f in files:
                    fp = os.path.join(root_dir, f)
                    if fp in _watched_files:
                        continue
                    try:
                        mtime = os.path.getmtime(fp)
                    except OSError:
                        continue
                    if mtime < start_time:
                        _watched_files.add(fp)
                        continue
                    _watched_files.add(fp)
                    ext = Path(f).suffix.lower()
                    if ext not in CONVERTIBLE_EXTS:
                        continue
                    stem = Path(f).stem
                    # 标记为转化中
                    _file_convert_status[stem] = {"status": "converting", "path": fp, "md_path": "", "name": f}
                    _new_files_pending.append(stem)
                    LOG.info(f"[文件监控] 检测到新文件: {f}，开始静默转化")
                    # 等待文件写入完成
                    time.sleep(2)
                    try:
                        tmp_path = FILE_SAVE_DIR / f
                        shutil.copy2(fp, tmp_path)
                        md_text, was_converted = convert_to_markdown(str(tmp_path))
                        if tmp_path.suffix.lower() not in (".md", ".txt"):
                            tmp_path.unlink(missing_ok=True)
                        if not md_text:
                            LOG.warning(f"[文件监控] 文件解析为空: {f}")
                            _file_convert_status[stem]["status"] = "error"
                            continue
                        save_path = FILE_SAVE_DIR / f"{stem}.md"
                        save_path.write_text(md_text, encoding="utf-8")
                        _file_convert_status[stem]["status"] = "done"
                        _file_convert_status[stem]["md_path"] = str(save_path)
                        LOG.info(f"[文件监控] 转化完成: {save_path}，{len(md_text)} 字符")
                    except Exception as e:
                        LOG.error(f"[文件监控] 转化失败 {f}: {e}", exc_info=True)
                        _file_convert_status[stem]["status"] = "error"
        except Exception as e:
            LOG.error(f"[文件监控] 扫描异常: {e}", exc_info=True)
            time.sleep(5)


def ask_ai_with_image(session_id: str, image_path: str, user_text: str, context_lines: list = None) -> str:
    if AI_PROVIDER not in VISION_PROVIDERS:
        return f"[当前模型 {AI_MODEL} 不支持图片，请切换到 gpt 或 claude]"
    try:
        with open(image_path, "rb") as f:
            img_b64 = base64.b64encode(f.read()).decode()
        ext = Path(image_path).suffix.lower().strip(".")
        mime = "jpeg" if ext in ("jpg", "jpeg") else ext

        if context_lines:
            context_block = "\n".join(context_lines)
            text_part = f"以下是群里最近的聊天记录供你参考：\n{context_block}\n\n现在有人发了一张图片并问：{user_text or '请描述这张图片'}"
        else:
            text_part = user_text or "请描述这张图片"

        messages = [{"role": "system", "content": SYSTEM_PROMPT}, {"role": "user", "content": [{"type": "text", "text": text_part}, {"type": "image_url", "image_url": {"url": f"data:image/{mime};base64,{img_b64}"}}]}]
        response = client.chat.completions.create(model=AI_MODEL, max_tokens=4096, messages=messages)
        reply = response.choices[0].message.content
        if session_id not in conversation_history:
            conversation_history[session_id] = []
        conversation_history[session_id].append({"role": "user", "content": text_part + " [附图]"})
        conversation_history[session_id].append({"role": "assistant", "content": reply})
        return reply
    except Exception as e:
        LOG.error(f"图片处理失败: {e}")
        return f"[图片处理失败: {e}]"


def ask_ai(session_id: str, user_text: str, context_lines: list = None) -> str:
    if session_id not in conversation_history:
        conversation_history[session_id] = []

    if context_lines:
        context_block = "\n".join(context_lines)
        full_query = f"以下是群里最近的聊天记录供你参考：\n{context_block}\n\n现在有人问你：{user_text}"
    else:
        full_query = user_text

    conversation_history[session_id].append({"role": "user", "content": full_query})

    if len(conversation_history[session_id]) > 40:
        conversation_history[session_id] = conversation_history[session_id][-40:]

    try:
        memory_block = memory_load_all()
        system_content = SYSTEM_PROMPT
        if memory_block:
            system_content += f"\n\n以下是你需要记住的长期记忆：\n{memory_block}"
        response = client.chat.completions.create(
            model=AI_MODEL,
            max_tokens=4096,
            messages=[{"role": "system", "content": system_content}] + conversation_history[session_id],
        )
        reply = response.choices[0].message.content
        conversation_history[session_id].append({"role": "assistant", "content": reply})
        return reply
    except Exception as e:
        LOG.error(f"AI API 调用失败: {e}")
        return f"[AI 暂时不可用: {e}]"


def handle_msg(wcf: Wcf, msg: WxMsg):
    try:
        content = msg.content or ""
        is_group = msg.from_group()
        self_wxid = wcf.get_self_wxid()

        LOG.debug(f"handle_msg: type={msg.type} is_group={is_group} roomid={msg.roomid} sender={msg.sender} is_at={msg.is_at(self_wxid)} content={content[:60]}")

        if msg.from_self():
            return

        # 过滤公众号/服务号，防止回复死循环
        if msg.sender.startswith("gh_") or msg.roomid.startswith("gh_"):
            return

        if is_group:
            if GROUP_WHITELIST and group_whitelist_ids:
                if msg.roomid not in group_whitelist_ids:
                    return

            is_at_me = msg.is_at(self_wxid)
            session_id = msg.roomid
            sender_name = wcf.get_alias_in_chatroom(msg.sender, msg.roomid) or msg.sender

            # 旁听：普通文字消息存缓冲区，顺便做时间关键词检测
            if msg.type == 1 and not is_at_me:
                buf = group_context_buffer.setdefault(msg.roomid, [])
                buf.append(f"{sender_name}: {content}")
                if len(buf) > GROUP_CONTEXT_LIMIT:
                    buf.pop(0)
                # 清理超时的待确认日程（5分钟）
                stale = [r for r, v in pending_schedule_confirm.items() if time.time() - v["ts"] > 300]
                for r in stale:
                    pending_schedule_confirm.pop(r, None)
                # 正则预筛：有时间关键词才调AI
                if _has_time_keywords(content):
                    def _async_detect(roomid, text, sname):
                        times = _extract_times_with_ai(text)
                        if not times:
                            return
                        pending_schedule_confirm[roomid] = {"content": text, "times": times, "ts": time.time(), "sender": sname}
                        wcf_ref = _global_wcf
                        if not wcf_ref:
                            return
                        if len(times) == 1:
                            t = times[0]
                            try:
                                dt = datetime.strptime(t["datetime"], "%Y-%m-%d %H:%M")
                                wd_str = f"周{WEEKDAYS_ZH[dt.weekday()]}"
                                hint = f"{dt.day:02d}/{dt.month:02d}/{str(dt.year)[2:]},{dt.hour:02d}:{dt.minute:02d},{wd_str}"
                            except Exception:
                                hint = t["datetime"]
                            wcf_ref.send_text(
                                f"🗓 检测到日程：{t['desc']}\n"
                                f"是否创建？请确认时间或忽略：\n@W. {hint}",
                                roomid
                            )
                        else:
                            lines = ["🗓 检测到多个时间点，请选择并确认："]
                            for t in times:
                                try:
                                    dt = datetime.strptime(t["datetime"], "%Y-%m-%d %H:%M")
                                    wd_str = f"周{WEEKDAYS_ZH[dt.weekday()]}"
                                    hint = f"{dt.day:02d}/{dt.month:02d}/{str(dt.year)[2:]},{dt.hour:02d}:{dt.minute:02d},{wd_str}"
                                except Exception:
                                    hint = t["datetime"]
                                lines.append(f"  · {t['desc']} → @W. {hint}")
                            wcf_ref = _global_wcf
                            if wcf_ref:
                                wcf_ref.send_text("\n".join(lines), roomid)
                    Thread(target=_async_detect, args=(msg.roomid, content, sender_name), daemon=True).start()
                return

            # 图片消息：存 pending；若未 @bot 则旁听返回，否则继续触发
            if msg.type == 3:
                group_pending_media[msg.roomid] = {"type": "image", "msg": msg}
                buf = group_context_buffer.setdefault(msg.roomid, [])
                buf.append(f"{sender_name}: [发送了一张图片]")
                if not is_at_me:
                    return

            # type=49(文件/链接)：DLL层会导致微信crash，这里不做处理
            # 文件靠 _file_watcher 自动检测 + 用户 "读取文件" 指令按需读取
            if msg.type == 49:
                return

            if not is_at_me:
                if GROUP_TRIGGER_PREFIX and not content.startswith(GROUP_TRIGGER_PREFIX):
                    return
                elif not GROUP_TRIGGER_PREFIX:
                    return

            # === 以下是 @bot 触发后的处理 ===
            query = content
            if GROUP_TRIGGER_PREFIX and content.startswith(GROUP_TRIGGER_PREFIX):
                query = content[len(GROUP_TRIGGER_PREFIX):].strip()
            elif is_at_me:
                try:
                    at_tag = f"@{wcf.get_user_info()['name']}"
                except Exception:
                    at_tag = ""
                if at_tag and at_tag in content:
                    query = content.replace(at_tag, "").strip()
                else:
                    # 备用：用regex去掉 @xxx 前缀
                    query = re.sub(r'^@\S+\s*', '', content).strip()
            LOG.info(f"[剥离] content={content[:60]!r} → query={query[:60]!r}")

            context_lines = group_context_buffer.pop(msg.roomid, [])
            pending = group_pending_media.pop(msg.roomid, None)
            LOG.info(f"群聊 [{msg.roomid}] 收到: {query[:60]}，上下文 {len(context_lines)} 条，pending={pending and pending['type']}")

            # 有待处理图片
            if pending and pending["type"] == "image":
                pm = pending["msg"]
                os.makedirs("D:\\Weixin\\bot\\downloads", exist_ok=True)
                img_path = wcf.download_image(pm.id, pm.extra, "D:\\Weixin\\bot\\downloads", 30)
                if img_path:
                    LOG.info(f"图片已下载: {img_path}")
                    reply = ask_ai_with_image(session_id, img_path, query, context_lines)
                else:
                    reply = "[图片下载失败，请重试]"
                wcf.send_text(reply, msg.roomid, msg.sender)
                return

        else:
            if not PRIVATE_CHAT_REPLY:
                return
            if msg.type not in (1, 3, 49):
                return
            context_lines = []
            session_id = msg.sender

            if msg.type == 3:
                time.sleep(1)
                img_path = wcf.download_image(msg.id, msg.extra, "D:\\Weixin\\bot\\downloads", 10)
                if img_path:
                    reply = ask_ai_with_image(session_id, img_path, "", [])
                else:
                    reply = "[图片下载失败，请重试]"
                wcf.send_text(reply, msg.sender)
                return

            query = content.strip()
            LOG.info(f"私聊 [{msg.sender}] 收到: {query[:60]}")

        if not query:
            HELP_TEXT = (
                "📋 可用指令：\n"
                "\n"
                "【文件】\n"
                "直接发文件即可，自动识别并转化\n"
                "列出文件 / ls文件\n"
                "总结/读取/分析/翻译文件 文件名[，问题]\n"
                "\n"
                "【笔记】\n"
                "新增笔记 笔记名，内容\n"
                "列出笔记\n"
                "查看笔记 笔记名\n"
                "\n"
                "【日程】\n"
                "创建日程 DD/MM/YY HH:MM 名称[，内容]\n"
                "列出日程\n"
                "查询日程 名称/日期/周几\n"
                "（群消息含时间时自动询问是否创建日程）\n"
                "\n"
                "【记忆】\n"
                "列出记忆 / ls\n"
                "记住 名字：内容\n"
                "读取记忆 名字\n"
                "删除记忆 名字\n"
                "导入文件为记忆 文件名[，记忆名]\n"
                "\n"
                "【其他】\n"
                "clear / 重置 → 清除对话上下文\n"
                "直接提问 → AI 回答\n"
                "单独 @bot → 显示此帮助"
            )
            wcf.send_text(HELP_TEXT, msg.roomid if is_group else msg.sender)
            return

        # === 记忆指令拦截 ===
        q = query.strip()

        # 日程确认：DD/MM/YY，HH:MM，周X
        m = re.match(r'^(\d{1,2})/(\d{1,2})/(\d{2,4})[，,](\d{2}):(\d{2})[，,]周([一二三四五六日天])$', q)
        if m and is_group and msg.roomid in pending_schedule_confirm:
            pending = pending_schedule_confirm.pop(msg.roomid)
            dd, mo, yy, hh, mi, wd_ch = m.group(1), m.group(2), m.group(3), m.group(4), m.group(5), m.group(6)
            year = int(yy) + 2000 if len(yy) <= 2 else int(yy)
            wd_map = {"一": 0, "二": 1, "三": 2, "四": 3, "五": 4, "六": 5, "日": 6, "天": 6}
            schedule_add(
                name=pending["times"][0]["desc"] if len(pending["times"]) == 1 else f"{dd}/{mo}/{yy} 日程",
                content=pending["content"],
                year=year, month=int(mo), day=int(dd),
                hour=int(hh), minute=int(mi),
                weekday=wd_map.get(wd_ch, 0),
                roomid=msg.roomid
            )
            wcf.send_text(
                f"✅ 日程已创建：{year}/{int(mo):02d}/{int(dd):02d} {hh}:{mi} 周{wd_ch}",
                msg.roomid, msg.sender
            )
            return

        # 清除上下文
        if q in ("clear", "清除上下文", "清空上下文", "重置"):
            conversation_history.pop(session_id, None)
            group_context_buffer.pop(msg.roomid, None) if is_group else None
            group_pending_media.pop(msg.roomid, None) if is_group else None
            wcf.send_text("上下文已清除。", msg.roomid if is_group else msg.sender)
            return

        # 列出记忆
        if q in ("列出记忆", "查看记忆", "ls", "ls memory"):
            wcf.send_text(memory_list(), msg.roomid if is_group else msg.sender)
            return

        # 删除记忆：删除记忆 xxx
        m = re.match(r'^删除记忆[：:\s]+(.+)$', q)
        if m:
            name = m.group(1).strip()
            if memory_delete(name):
                wcf.send_text(f"记忆 {name}.txt 已删除。", msg.roomid if is_group else msg.sender)
            else:
                wcf.send_text(f"未找到记忆 {name}.txt。", msg.roomid if is_group else msg.sender)
            return

        # 读取记忆：读取记忆 xxx
        m = re.match(r'^读取记忆[：:\s]+(.+)$', q)
        if m:
            name = m.group(1).strip()
            content = memory_read(name)
            if content:
                wcf.send_text(f"📄 {name}.txt：\n{content}", msg.roomid if is_group else msg.sender)
            else:
                wcf.send_text(f"未找到记忆 {name}.txt。", msg.roomid if is_group else msg.sender)
            return

        # 导入文件为记忆：导入文件为记忆 文件名[，记忆名]
        m = re.match(r'^导入文件为记忆[\s：:]+(.+)$', q)
        if m:
            raw = m.group(1).strip()
            parts = re.split(r'[，,]\s*', raw, maxsplit=1)
            src_name = parts[0].strip()
            mem_name = parts[1].strip() if len(parts) > 1 else Path(src_name).stem
            # 搜索顺序：FILE_SAVE_DIR → 微信缓存
            found_path = None
            for f in FILE_SAVE_DIR.iterdir():
                if f.stem == src_name or f.name == src_name or f.stem.startswith(src_name) or f.name.startswith(src_name):
                    found_path = str(f)
                    break
            if not found_path:
                for root_dir, dirs, files in os.walk(WECHAT_FILE_DIR):
                    for f in files:
                        if f == src_name or f.startswith(src_name):
                            found_path = os.path.join(root_dir, f)
                            break
                    if found_path:
                        break
            if not found_path:
                wcf.send_text(f"[未找到文件《{src_name}》]", msg.roomid if is_group else msg.sender)
                return
            file_text = extract_file_text(found_path)
            if not file_text:
                wcf.send_text(f"[文件内容为空或解析失败]", msg.roomid if is_group else msg.sender)
                return
            memory_save(mem_name, file_text)
            wcf.send_text(f"✅ 已将《{Path(found_path).name}》保存为记忆「{mem_name}」（{len(file_text)} 字符）", msg.roomid if is_group else msg.sender)
            return

        # ── 笔记指令 ─────────────────────────────────────────────────────
        if q in ("列出笔记", "ls笔记"):
            wcf.send_text(note_list(), msg.roomid if is_group else msg.sender)
            return

        m = re.match(r'^新增笔记[\s：:]+(.+)$', q, re.DOTALL)
        if m:
            raw = m.group(1).strip()
            parts = re.split(r'[，,]', raw, maxsplit=1)
            nname = parts[0].strip()
            ncontent = parts[1].strip() if len(parts) > 1 else ""
            if not ncontent:
                wcf.send_text("请提供笔记内容，格式：新增笔记 笔记名，内容", msg.roomid if is_group else msg.sender)
            else:
                note_save(nname, ncontent)
                wcf.send_text(f"📝 笔记「{nname}」已保存。", msg.roomid if is_group else msg.sender)
            return

        m = re.match(r'^查看笔记[\s：:]+(.+)$', q)
        if m:
            text, stem = note_read(m.group(1).strip())
            if text is None:
                wcf.send_text(f"未找到笔记「{m.group(1).strip()}」。", msg.roomid if is_group else msg.sender)
            else:
                wcf.send_text(f"📝 {stem}：\n{text}", msg.roomid if is_group else msg.sender)
            return

        # ── 日程指令 ─────────────────────────────────────────────────────
        if q in ("列出日程", "ls日程"):
            wcf.send_text(schedule_list(), msg.roomid if is_group else msg.sender)
            return

        m = re.match(r'^查询日程[\s：:]+(.+)$', q)
        if m:
            wcf.send_text(schedule_query(m.group(1).strip()), msg.roomid if is_group else msg.sender)
            return

        m = re.match(r'^创建日程(.*)$', q)
        if m:
            rest = m.group(1).strip()
            reply_to = msg.roomid if is_group else msg.sender
            TMPL = (
                "📅 日程创建格式：\n"
                "创建日程 DD/MM/YY HH:MM 名称[，内容]\n\n"
                "示例：\n"
                "创建日程 30/03/26 14:00 课题汇报，讨论钛合金分子动力学问题"
            )
            if not rest:
                wcf.send_text(TMPL, reply_to)
                return
            pm = re.match(r'^(\d{1,2})/(\d{1,2})/(\d{2,4})\s+(\d{2}):(\d{2})\s+([^，,]+)[，,]?(.*)$', rest)
            if not pm:
                wcf.send_text(f"格式不对，请参考：\n{TMPL}", reply_to)
                return
            dd, mo, yy, hh, mi, name, content_part = pm.group(1,2,3,4,5,6,7)
            year = int(yy) + 2000 if len(yy) <= 2 else int(yy)
            dt = datetime(year, int(mo), int(dd), int(hh), int(mi))
            schedule_add(
                name=name.strip(), content=content_part.strip(),
                year=year, month=int(mo), day=int(dd),
                hour=int(hh), minute=int(mi),
                weekday=dt.weekday(),
                roomid=msg.roomid if is_group else msg.sender
            )
            wcf.send_text(
                f"✅ 日程已创建：{name.strip()}\n{year}/{int(mo):02d}/{int(dd):02d} {hh}:{mi} 周{WEEKDAYS_ZH[dt.weekday()]}",
                reply_to
            )
            return

        # 列出文件：列出文件 / ls文件 / 文件列表
        if q in ("列出文件", "ls文件", "文件列表", "ls files"):
            reply_to = msg.roomid if is_group else msg.sender
            seen_stems = set()
            lines = []
            status_map = {"converting": "⏳转化中", "done": "✅", "error": "❌"}
            # 1) 转化状态中的文件
            for stem, info in _file_convert_status.items():
                s = status_map.get(info.get("status", ""), "❓")
                lines.append(f"  {s} {info.get('name', stem)}")
                seen_stems.add(stem)
            # 2) FILE_SAVE_DIR 中已保存的 md 文件
            for f in FILE_SAVE_DIR.iterdir():
                if f.suffix == ".md" and f.stem not in seen_stems:
                    lines.append(f"  ✅ {f.stem}")
                    seen_stems.add(f.stem)
            # 3) WECHAT_FILE_DIR 中可转化但尚未处理的文件（最近20个）
            raw_files = []
            for root_dir, dirs, files in os.walk(WECHAT_FILE_DIR):
                for f in files:
                    p = Path(f)
                    if p.suffix.lower() in CONVERTIBLE_EXTS and p.stem not in seen_stems:
                        fp = os.path.join(root_dir, f)
                        raw_files.append((f, os.path.getmtime(fp)))
            raw_files.sort(key=lambda x: x[1], reverse=True)
            for f, _ in raw_files[:20]:
                lines.append(f"  📄 {f}")
                seen_stems.add(Path(f).stem)
            # 4) DB中最近1小时的文件消息（尚未下载的）
            try:
                min_ts = int(time.time()) - 3600
                msg_dbs = [d for d in wcf.get_dbs() if d.startswith("MSG")]
                for db in msg_dbs:
                    tables = wcf.get_tables(db)
                    tbl_names = [t.get("name", "") for t in tables] if isinstance(tables, list) else []
                    for tbl in tbl_names:
                        if not tbl.startswith("MSG"):
                            continue
                        try:
                            rows = wcf.query_sql(db,
                                f"SELECT StrContent FROM {tbl} WHERE Type=49 AND CreateTime>{min_ts} ORDER BY CreateTime DESC LIMIT 20")
                            for row in rows:
                                raw = row.get("StrContent", "")
                                if isinstance(raw, bytes):
                                    try:
                                        raw = raw.decode("utf-8")
                                    except Exception:
                                        continue
                                fname = _extract_filename_from_xml(raw)
                                if fname and Path(fname).suffix.lower() in CONVERTIBLE_EXTS and Path(fname).stem not in seen_stems:
                                    lines.append(f"  ⬇️ {fname}")
                                    seen_stems.add(Path(fname).stem)
                        except Exception:
                            pass
            except Exception as e:
                LOG.debug(f"[ls文件] DB查询失败: {e}")
            if lines:
                wcf.send_text(f"📂 文件列表（{len(lines)} 个）：\n" + "\n".join(lines) + "\n\n✅=已转化  📄=本地缓存  ⬇️=可从微信下载\n用法：读取文件 文件名", reply_to)
            else:
                wcf.send_text("[当前无已检测或已保存的文件]", reply_to)
            return

        # 读取/总结文件：读取文件 xxx 或 总结文件 xxx[，问题]
        m = re.match(r'^(读取|总结|分析|翻译)文件[\s：:]+(.+)$', q)
        if m:
            action = m.group(1).strip()
            raw = m.group(2).strip()
            reply_to = msg.roomid if is_group else msg.sender
            # 逗号后面的内容作为附加问题，前面才是文件名
            parts = re.split(r'[，,]+', raw, maxsplit=1)
            filename = parts[0].strip()
            extra_q = parts[1].strip() if len(parts) > 1 else ""
            LOG.info(f"[读取文件] 搜索: {filename!r}")

            # ---------- 收集所有匹配 ----------
            candidates = []  # [(display_name, full_path, mtime, source)]
            seen_paths = set()

            def _add(path_str, source):
                rp = os.path.normpath(path_str)
                if rp in seen_paths:
                    return
                seen_paths.add(rp)
                try:
                    mt = os.path.getmtime(rp)
                except OSError:
                    mt = 0
                candidates.append((Path(rp).name, rp, mt, source))

            # ① FILE_SAVE_DIR（已转化的 .md/.txt 才算）
            cached_stems = set()
            if FILE_SAVE_DIR.exists():
                for f in FILE_SAVE_DIR.iterdir():
                    if f.suffix.lower() not in (".md", ".txt"):
                        continue
                    if f.stem == filename or f.name == filename or f.stem.startswith(filename):
                        _add(str(f), "已转化")
                        cached_stems.add(f.stem)

            # ② WECHAT_FILE_DIR（微信本地缓存，跳过已有缓存MD的）
            for root_dir, _dirs, files in os.walk(WECHAT_FILE_DIR):
                for f in files:
                    if f == filename or f.startswith(filename) or Path(f).stem == filename:
                        fp = os.path.join(root_dir, f)
                        ext_l = Path(f).suffix.lower()
                        if ext_l in CONVERTIBLE_EXTS and Path(f).stem not in cached_stems:
                            _add(fp, "本地缓存")

            # ③ DB 尝试（失败不抛异常）
            if not candidates:
                try:
                    db_path = _try_download_from_db(wcf, filename)
                    if db_path:
                        _add(db_path, "微信下载")
                except Exception as e:
                    LOG.warning(f"[读取文件] DB下载尝试失败: {e}")

            if not candidates:
                wcf.send_text(f"[未找到文件《{filename}》，请确认文件名或重新发送]", reply_to)
                return

            # ---------- 多匹配：列出让用户选 ----------
            if len(candidates) > 1:
                lines = []
                for name, path, mt, src in sorted(candidates, key=lambda x: x[2], reverse=True):
                    ts = time.strftime("%m-%d %H:%M", time.localtime(mt)) if mt else "未知"
                    lines.append(f"  [{src}] {name}  ({ts})")
                wcf.send_text(
                    f"🔍 找到 {len(candidates)} 个匹配《{filename}》的文件：\n"
                    + "\n".join(lines)
                    + "\n\n请用更精确的文件名重试，如：\n读取文件 完整文件名.docx",
                    reply_to
                )
                return

            # ---------- 单匹配：读取 ----------
            chosen_name, chosen_path, _, chosen_src = candidates[0]
            ext = Path(chosen_path).suffix.lower()
            LOG.info(f"[读取文件] 命中: {chosen_path} ({chosen_src})")

            # 如果已是 .md/.txt，直接读
            if ext in (".md", ".txt"):
                file_text = Path(chosen_path).read_text(encoding="utf-8", errors="replace")
            else:
                # 检查是否已有缓存的 .md
                cached_md = FILE_SAVE_DIR / f"{Path(chosen_path).stem}.md"
                if cached_md.exists():
                    file_text = cached_md.read_text(encoding="utf-8", errors="replace")
                    LOG.info(f"[读取文件] 使用缓存MD: {cached_md}")
                else:
                    # 转化并缓存
                    tmp_path = FILE_SAVE_DIR / Path(chosen_path).name
                    convert_src = chosen_path  # 默认从原路径转化
                    if tmp_path != Path(chosen_path):
                        try:
                            if tmp_path.exists():
                                tmp_path.unlink()
                            shutil.copy2(chosen_path, tmp_path)
                            convert_src = str(tmp_path)
                        except Exception as e:
                            LOG.warning(f"[读取文件] 复制失败，直接从原路径转化: {e}")
                            convert_src = chosen_path
                            tmp_path = None
                    file_text, _ = convert_to_markdown(convert_src)
                    if tmp_path and tmp_path.suffix.lower() not in (".md", ".txt"):
                        try:
                            tmp_path.unlink(missing_ok=True)
                        except PermissionError:
                            LOG.debug(f"[读取文件] 临时文件删除失败(文件锁): {tmp_path}")
                    if file_text:
                        cached_md.write_text(file_text, encoding="utf-8")
                        LOG.info(f"[读取文件] 转化并缓存: {cached_md} ({len(file_text)} 字符)")
                    else:
                        wcf.send_text(f"[文件《{chosen_name}》解析失败，格式可能不支持]", reply_to)
                        return

            if not file_text or not file_text.strip():
                wcf.send_text(f"[文件《{chosen_name}》内容为空]", reply_to)
                return
            LOG.info(f"[读取文件] {chosen_name}，{len(file_text)} 字符，发给AI")
            task = extra_q if extra_q else f"请{action}这个文件的内容"
            query = (
                f"用户上传了文件《{chosen_name}》，以下是该文件的完整文本内容（已由系统提取）。"
                f"请直接基于这些内容回答，不要说你无法读取文件。\n\n"
                f"──── 文件内容 ────\n{file_text[:100000]}\n──── 文件结束 ────\n\n"
                f"用户的要求：{task}"
            )

        # 保存记忆：记住 xxx（让AI生成内容保存）或 记住 名字：内容
        m = re.match(r'^记住[：::\s]+([^：:]+)[：:](.+)$', q, re.DOTALL)
        if m:
            name = m.group(1).strip()
            content = m.group(2).strip()
            saved = memory_save(name, content)
            wcf.send_text(f"相关记忆已保存到 {saved}。", msg.roomid if is_group else msg.sender)
            return

        time.sleep(0.5)
        reply = ask_ai(session_id, query, context_lines if is_group else None)
        LOG.info(f"回复 [{session_id}]: {reply[:80]}")

        if is_group:
            wcf.send_text(reply, msg.roomid, msg.sender)
        else:
            wcf.send_text(reply, msg.sender)

    except Exception as e:
        LOG.error(f"处理消息出错: {e}", exc_info=True)


def safe_handle(wcf: Wcf, msg: WxMsg):
    try:
        handle_msg(wcf, msg)
    except Exception as e:
        LOG.error(f"handle_msg 未捕获异常: {e}", exc_info=True)


def recv_loop(wcf: Wcf):
    LOG.info("消息接收循环已启动")
    empty_count = 0
    HEARTBEAT_INTERVAL = 10  # 连续10次Empty后检测心跳
    while wcf.is_receiving_msg():
        try:
            msg = wcf.get_msg()
            empty_count = 0  # 收到消息，重置计数
            LOG.debug(f"收到原始消息 type={msg.type} roomid={msg.roomid} sender={msg.sender}")
            if msg.type == 49:
                LOG.debug("type=49 跳过（DLL层crash风险）")
                continue
            Thread(target=safe_handle, args=(wcf, msg), daemon=True).start()
        except Empty:
            empty_count += 1
            if empty_count >= HEARTBEAT_INTERVAL:
                empty_count = 0
                # 心跳：检查 WeChat.exe 进程是否存活（比 wcf API 快得多）
                result = subprocess.run(
                    ["tasklist", "/FI", "IMAGENAME eq WeChat.exe", "/NH"],
                    capture_output=True, text=True, timeout=5
                )
                if "WeChat.exe" not in result.stdout:
                    LOG.warning("心跳检测：WeChat.exe 进程不存在，微信已退出")
                    return
        except Exception as e:
            LOG.error(f"接收消息出错: {e}", exc_info=True)


def _auto_click_login():
    """后台线程：监控微信登录窗口自动点击登录，同时监控 WxInitSDK 注入失败弹窗自动关闭。"""
    try:
        from pywinauto import Application, Desktop
    except ImportError:
        LOG.warning("[自动登录] pywinauto 未安装，跳过")
        return
    LOG.info("[自动登录] 后台监控已启动，等待登录窗口...")
    login_clicked = False
    for _ in range(60):
        try:
            # 优先检测"注入失败"弹窗并自动关闭
            try:
                for win in Desktop(backend='uia').windows():
                    if 'WxInitSDK' in win.window_text() or '注入失败' in win.window_text():
                        try:
                            btn = win.child_window(title='确定', control_type='Button')
                            if btn.exists(timeout=0):
                                btn.click_input()
                                LOG.warning("[自动登录] 已关闭「注入失败」弹窗")
                        except Exception:
                            try:
                                win.close()
                            except Exception:
                                pass
            except Exception:
                pass

            if login_clicked:
                time.sleep(1)
                continue

            # 检测微信登录窗口
            app = Application(backend='uia').connect(path='WeChat.exe', timeout=2)
            dlg = app.window(title='微信')
            if not dlg.exists(timeout=1):
                time.sleep(1)
                continue
            for btn_name in ['进入微信', '登录', '登錄', 'Enter WeChat', 'Log In']:
                try:
                    btn = dlg.child_window(title=btn_name, control_type='Button')
                    if btn.exists(timeout=1):
                        btn.click_input()
                        LOG.info(f"[自动登录] 已点击「{btn_name}」按钮")
                        login_clicked = True
                        break
                except Exception:
                    continue
            if not login_clicked:
                try:
                    dlg.set_focus()
                    import pywinauto.keyboard as kb
                    kb.send_keys('{ENTER}')
                    LOG.info("[自动登录] 已发送 Enter 键")
                    login_clicked = True
                except Exception:
                    pass
        except Exception:
            time.sleep(1)
    LOG.info("[自动登录] 监控线程结束")


def _kill_wechat():
    """强制结束所有 WeChat 相关进程 + 释放 wcferry 端口，确保注入前环境干净"""
    # 第一轮：wmic 删除所有 WeChat 相关进程（比 taskkill 更彻底）
    try:
        subprocess.run(
            ['wmic', 'process', 'where', "name like '%WeChat%'", 'delete'],
            capture_output=True, timeout=10
        )
    except Exception:
        pass
    # 第二轮：taskkill 兜底
    for proc_name in ["WeChat.exe", "WeChatUtility.exe", "WeChatPlayer.exe",
                       "WeChatBrowser.exe", "WeChatAppEx.exe"]:
        try:
            subprocess.run(["taskkill", "/F", "/IM", proc_name, "/T"],
                           capture_output=True, timeout=5)
        except Exception:
            pass
    # 释放 wcferry 占用的端口（10086/10087）
    for port in [10086, 10087]:
        try:
            r = subprocess.run(
                ["netstat", "-ano", "-p", "TCP"],
                capture_output=True, text=True, timeout=5
            )
            for line in r.stdout.splitlines():
                if f":{port}" in line and ("LISTENING" in line or "ESTABLISHED" in line):
                    parts = line.split()
                    pid = parts[-1]
                    if pid.isdigit() and int(pid) > 0:
                        subprocess.run(["taskkill", "/F", "/PID", pid],
                                       capture_output=True, timeout=5)
                        LOG.info(f"[清理] 已杀死占用端口 {port} 的进程 PID={pid}")
        except Exception:
            pass
    # 等待进程完全退出
    time.sleep(5)
    # 验证是否清理干净
    r = subprocess.run(
        ["tasklist", "/FI", "IMAGENAME eq WeChat.exe", "/NH"],
        capture_output=True, text=True, timeout=5
    )
    if "WeChat.exe" in r.stdout:
        LOG.warning("[清理] WeChat.exe 仍在运行，再次强杀...")
        subprocess.run(["taskkill", "/F", "/IM", "WeChat.exe"], capture_output=True, timeout=5)
        time.sleep(3)
    LOG.info("[清理] 微信进程和端口清理完成")


class _WcfInitError(Exception):
    """Wcf() 初始化失败（注入失败或连接失败）时抛出，防止 os._exit 杀死进程"""
    pass


def _init_wcf():
    """初始化 wcferry 连接，解析白名单群，返回 Wcf 实例。
    Monkey-patch os._exit 防止 wcferry 注入失败时杀死整个进程，改为可重试的异常。"""
    global _global_wcf
    _real_exit = os._exit

    def _fake_exit(code):
        raise _WcfInitError(f"wcferry 调用了 os._exit({code})，已拦截")

    os._exit = _fake_exit
    try:
        wcf = Wcf()
    finally:
        os._exit = _real_exit  # 无论成功失败都恢复

    wcf.enable_receiving_msg()
    _global_wcf = wcf
    LOG.info(f"微信连接成功，当前账号 wxid: {wcf.get_self_wxid()}")
    info = wcf.get_user_info()
    LOG.info(f"账号昵称: {info.get('name')}  微信号: {info.get('wxid')}")

    if GROUP_WHITELIST:
        group_whitelist_ids.clear()
        contacts = wcf.get_contacts()
        for c in contacts:
            if c.get('name') in GROUP_WHITELIST:
                group_whitelist_ids.add(c.get('wxid'))
                LOG.info(f"白名单群: {c.get('name')} -> {c.get('wxid')}")
        if not group_whitelist_ids:
            LOG.warning(f"未找到白名单群 {GROUP_WHITELIST}，将监听所有群（请确认群名正确）")
    else:
        LOG.info("未设置 GROUP_WHITELIST，将响应所有群的触发词消息")
    return wcf


def main():
    global _global_wcf, _file_watcher_started, _schedule_reminder_started
    LOG.info(f"正在启动 wcferry，AI 供应商: {AI_PROVIDER} ({AI_MODEL})")
    LOG.info("请确保微信 3.9.x 已登录...")

    _init_schedules_db()
    if not _schedule_reminder_started:
        _schedule_reminder_started = True
        Thread(target=_schedule_reminder_loop, daemon=True).start()
        LOG.info("[日程] 提醒线程已启动")

    global _push_server_started, _push_worker_started
    _ensure_push_queue()
    if not _push_server_started:
        _push_server_started = True
        Thread(target=_push_server, daemon=True).start()
    if not _push_worker_started:
        _push_worker_started = True
        Thread(target=_push_send_worker, daemon=True).start()

    RECONNECT_INTERVAL = 10  # 断开后每隔10秒尝试重连
    WECHAT_EXE = os.getenv("WECHAT_EXE", r"D:\vx3.9\WeChat\WeChat.exe")
    MAX_INJECT_RETRIES = 3   # 注入最多尝试次数

    while True:
        wcf = None
        try:
            for attempt in range(1, MAX_INJECT_RETRIES + 1):
                # 检查微信是否在运行
                r = subprocess.run(
                    ["tasklist", "/FI", "IMAGENAME eq WeChat.exe", "/NH"],
                    capture_output=True, text=True, timeout=5
                )
                wechat_running = "WeChat.exe" in r.stdout

                if not wechat_running:
                    if attempt > 1:
                        _kill_wechat()  # 清理残留
                    LOG.info(f"[注入] 微信未运行，手动启动: {WECHAT_EXE}")
                    subprocess.Popen([WECHAT_EXE], shell=False)
                    Thread(target=_auto_click_login, daemon=True).start()
                    LOG.info("[注入] 等待微信启动和登录（30秒）...")
                    time.sleep(30)
                else:
                    LOG.info("[注入] 微信已在运行")
                    Thread(target=_auto_click_login, daemon=True).start()

                try:
                    LOG.info(f"[注入] 第 {attempt}/{MAX_INJECT_RETRIES} 次尝试...")
                    wcf = _init_wcf()
                    break  # 注入成功
                except _WcfInitError as e:
                    LOG.warning(f"[注入] 第 {attempt} 次失败: {e}")
                    _kill_wechat()  # 失败后清理，为下一次做准备
                    if attempt >= MAX_INJECT_RETRIES:
                        raise

            # 文件监控线程只启动一次（使用 _global_wcf，重连后自动跟随）
            if not _file_watcher_started:
                _file_watcher_started = True
                Thread(target=_file_watcher, daemon=True).start()

            # 重连后：查微信数据库恢复 crash 前的文件消息
            Thread(target=_recover_files_from_db, args=(wcf, group_whitelist_ids), daemon=True).start()

            recv_loop(wcf)
            # recv_loop 正常退出 = wcferry 断开（微信crash或退出）
            LOG.warning("⚠️ 消息接收循环退出，微信可能已断开")

        except KeyboardInterrupt:
            LOG.info("收到 Ctrl+C，正在退出...")
            break
        except Exception as e:
            LOG.error(f"wcferry 连接异常: {e}", exc_info=True)
        finally:
            _global_wcf = None
            if wcf:
                try:
                    wcf.disable_receiving_msg()
                except Exception:
                    pass
                try:
                    wcf.cleanup()
                except Exception:
                    pass

        LOG.info(f"⏳ {RECONNECT_INTERVAL} 秒后尝试重新连接微信...")
        time.sleep(RECONNECT_INTERVAL)


# ── HTTP 推送服务 ───────────────────────────────────────────────────────
# Sub Recorder 等外部服务通过 POST /notify 发送推送请求
# 请求体: {"to": "wxid_xxx 或群id", "msg": "消息内容", "token": "...(optional)"}
# PUSH_PORT: 监听端口，默认 5700
# PUSH_TOKEN: 鉴权 token，留空则不鉴权（仅本机访问时可留空）

PUSH_PORT  = int(os.getenv("PUSH_PORT", "5700"))
PUSH_TOKEN = os.getenv("PUSH_TOKEN", "").strip()
_push_server_started  = False
_push_worker_started  = False
_push_queue = None  # 延迟初始化
_bot_start_time = time.time()
_ws_log_clients: list = []  # 保留兼容（已废弃）
_sse_log_clients: list = []  # SSE 日志订阅者（每个元素是 queue.Queue）


def _ensure_push_queue():
    global _push_queue
    if _push_queue is None:
        _push_queue = queue.Queue(maxsize=200)
    return _push_queue


def _push_send_worker():
    """单独线程消费推送队列，微信断开时暂停等待而不丢弃消息"""
    q = _ensure_push_queue()
    LOG.info("[Push] 发送 worker 已启动")
    while True:
        try:
            to, msg = q.get(timeout=2)
        except queue.Empty:
            continue
        while True:
            wcf = _global_wcf
            if wcf is not None:
                try:
                    wcf.send_text(msg, to)
                    LOG.info(f"[Push] 已发送到 {to}: {msg[:40]}...")
                    q.task_done()
                    break
                except Exception as e:
                    LOG.error(f"[Push] send_text 失败: {e}")
                    time.sleep(3)
            else:
                LOG.debug("[Push] 微信未连接，等待 3s 后重试...")
                time.sleep(3)


# ── SSE 日志广播 Handler ─────────────────────────────────────────
class _SseLogHandler(logging.Handler):
    """把日志推送到所有 SSE 客户端的队列（线程安全）"""
    def emit(self, record):
        if not _sse_log_clients:
            return
        try:
            entry = json.dumps({
                "t": self.formatTime(record, "%H:%M:%S"),
                "lvl": record.levelname,
                "msg": self.format(record),
            }, ensure_ascii=False)
        except Exception:
            return
        dead = []
        for q in list(_sse_log_clients):
            try:
                q.put_nowait(entry)
            except Exception:
                dead.append(q)
        for q in dead:
            try:
                _sse_log_clients.remove(q)
            except ValueError:
                pass


class _WsLogHandler(logging.Handler):  # 保留兼容，已不使用
    def emit(self, record):
        pass


def _push_server():
    """Flask Web 管理界面 + HTTP 推送服务"""
    from flask import Flask, request, jsonify, send_from_directory, Response, stream_with_context

    q = _ensure_push_queue()
    app = Flask(__name__, static_folder=None)

    import logging as _logging
    _sse_handler = _SseLogHandler()
    _sse_handler.setFormatter(logging.Formatter("%(message)s"))
    _logging.getLogger().addHandler(_sse_handler)

    webui_dir = Path(__file__).parent / "webui"

    # ── 静态文件（不拦截 api/ 路径）───────────────────────────────────────
    @app.route("/")
    def index():
        return send_from_directory(str(webui_dir), "index.html")

    @app.route("/<path:filename>")
    def static_files(filename):
        if filename.startswith(("api/", "ws/")):
            return jsonify({"error": "not found"}), 404
        return send_from_directory(str(webui_dir), filename)

    # ── 兼容旧接口 ──────────────────────────────────────────────────────
    @app.route("/health")
    def health():
        status = "connected" if _global_wcf else "disconnected"
        return jsonify({"ok": True, "wechat": status, "queue": q.qsize()})

    @app.route("/notify", methods=["POST"])
    def notify():
        body = request.get_json(silent=True) or {}
        if PUSH_TOKEN and body.get("token", "") != PUSH_TOKEN:
            return jsonify({"ok": False, "msg": "unauthorized"}), 401
        to  = str(body.get("to",  "")).strip()
        msg = str(body.get("msg", "")).strip()
        if not to or not msg:
            return jsonify({"ok": False, "msg": "missing to or msg"}), 400
        if q.full():
            return jsonify({"ok": False, "msg": "queue full"}), 503
        q.put((to, msg))
        LOG.info(f"[Push] 收到推送请求 -> {to}，队列长度: {q.qsize()}")
        return jsonify({"ok": True, "queued": q.qsize()})

    # ── 管理 API ────────────────────────────────────────────────────────
    @app.route("/api/status")
    def api_status():
        uptime = int(time.time() - _bot_start_time)
        h, r = divmod(uptime, 3600)
        m, s = divmod(r, 60)
        return jsonify({
            "wechat": "connected" if _global_wcf else "disconnected",
            "ai_provider": AI_PROVIDER,
            "ai_model": AI_MODEL,
            "uptime": f"{h:02d}:{m:02d}:{s:02d}",
            "queue": q.qsize(),
        })

    @app.route("/api/reconnect", methods=["POST"])
    def api_reconnect():
        wcf = _global_wcf
        if wcf:
            try:
                wcf.cleanup()
            except Exception:
                pass
        return jsonify({"ok": True, "msg": "重连信号已发送，请等待 bot 自动重连"})

    @app.route("/api/send", methods=["POST"])
    def api_send():
        body = request.get_json(silent=True) or {}
        to  = str(body.get("to",  "")).strip()
        msg = str(body.get("msg", "")).strip()
        if not to or not msg:
            return jsonify({"ok": False, "msg": "missing to or msg"}), 400
        if q.full():
            return jsonify({"ok": False, "msg": "queue full"}), 503
        q.put((to, msg))
        return jsonify({"ok": True, "queued": q.qsize()})

    @app.route("/api/schedules")
    def api_schedules():
        try:
            conn = sqlite3.connect(str(SCHEDULES_DB))
            rows = conn.execute(
                "SELECT id, name, year, month, day, hour, minute, weekday, content, reminded_morning, reminded_before, created_at "
                "FROM schedules ORDER BY year,month,day,hour,minute"
            ).fetchall()
            conn.close()
            result = []
            for r in rows:
                result.append({
                    "id": r[0], "name": r[1],
                    "datetime": f"{r[2]}/{r[3]:02d}/{r[4]:02d} {r[5]:02d}:{r[6]:02d}",
                    "weekday": f"周{WEEKDAYS_ZH[r[7]]}" if r[7] is not None else "",
                    "content": r[8] or "",
                    "reminded_morning": bool(r[9]), "reminded_before": bool(r[10]),
                    "created_at": r[11],
                })
            return jsonify(result)
        except Exception as e:
            return jsonify({"error": str(e)}), 500

    @app.route("/api/schedules/<int:sid>", methods=["DELETE"])
    def api_delete_schedule(sid):
        try:
            conn = sqlite3.connect(str(SCHEDULES_DB))
            conn.execute("DELETE FROM schedules WHERE id=?", (sid,))
            conn.commit()
            conn.close()
            return jsonify({"ok": True})
        except Exception as e:
            return jsonify({"error": str(e)}), 500

    @app.route("/api/notes")
    def api_notes():
        files = sorted(NOTES_DIR.glob("*.txt"), key=lambda p: p.stat().st_mtime, reverse=True)
        return jsonify([{"name": p.stem, "mtime": int(p.stat().st_mtime)} for p in files])

    @app.route("/api/notes/<name>")
    def api_note_read(name):
        text, stem = note_read(name)
        if text is None:
            return jsonify({"error": "not found"}), 404
        return jsonify({"name": stem, "content": text})

    @app.route("/api/notes/<name>", methods=["DELETE"])
    def api_delete_note(name):
        matches = [p for p in NOTES_DIR.glob("*.txt") if p.stem.lower() == name.lower()]
        if not matches:
            return jsonify({"error": "not found"}), 404
        matches[0].unlink()
        return jsonify({"ok": True})

    # ── SSE 日志流（替代 WebSocket，线程安全）──────────────────────────
    @app.route("/api/logs/stream")
    def api_logs_stream():
        client_q = queue.Queue(maxsize=500)
        _sse_log_clients.append(client_q)

        def generate():
            try:
                while True:
                    try:
                        msg = client_q.get(timeout=25)
                        yield f"data: {msg}\n\n"
                    except queue.Empty:
                        yield ": keepalive\n\n"
            finally:
                try:
                    _sse_log_clients.remove(client_q)
                except ValueError:
                    pass

        return Response(
            stream_with_context(generate()),
            mimetype="text/event-stream",
            headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"},
        )

    # ── 账号信息 ────────────────────────────────────────────────────────
    @app.route("/api/account")
    def api_account():
        wcf = _global_wcf
        if not wcf:
            return jsonify({"wxid": "", "name": "", "avatar": ""})
        try:
            info = wcf.get_user_info()
            if not isinstance(info, dict):
                info = {}
            wxid = info.get("wxid") or info.get("Wxid") or ""
            name = info.get("name") or info.get("Name") or info.get("nickName") or ""
            # 尝试多个可能的头像字段名
            avatar_url = (info.get("small_head_img_url") or info.get("headImgUrl")
                          or info.get("avatar") or info.get("head_img") or "")
            LOG.debug(f"[account] get_user_info keys: {list(info.keys())}")
        except Exception as e:
            LOG.warning(f"[account] get_user_info 失败: {e}")
            wxid, name, avatar_url = "", "", ""
        return jsonify({"wxid": wxid, "name": name, "avatar": avatar_url})

    # ── 系统资源 ─────────────────────────────────────────────────────────
    @app.route("/api/sysinfo")
    def api_sysinfo():
        try:
            import psutil as _ps
            cpu = _ps.cpu_percent(interval=0.2)
            mem = _ps.virtual_memory()
            return jsonify({"cpu": round(cpu, 1), "mem_used": mem.used, "mem_total": mem.total, "mem_pct": round(mem.percent, 1)})
        except Exception:
            return jsonify({"cpu": None, "mem_used": None, "mem_total": None, "mem_pct": None})

    # ── 文件管理 ─────────────────────────────────────────────────────────
    @app.route("/api/files/raw")
    def api_files_raw():
        exts = {".docx", ".doc", ".xlsx", ".xls", ".pptx", ".ppt", ".pdf", ".txt", ".md"}
        try:
            files = [p for p in Path(WECHAT_FILE_DIR).iterdir() if p.suffix.lower() in exts]
            files.sort(key=lambda p: p.stat().st_mtime, reverse=True)
            return jsonify([{"name": p.name, "size": p.stat().st_size, "mtime": int(p.stat().st_mtime)} for p in files])
        except Exception as e:
            return jsonify({"error": str(e)}), 500

    @app.route("/api/files/converted")
    def api_files_converted():
        try:
            files = sorted(FILE_SAVE_DIR.glob("*"), key=lambda p: p.stat().st_mtime, reverse=True)
            return jsonify([{"name": p.name, "size": p.stat().st_size, "mtime": int(p.stat().st_mtime)} for p in files if p.is_file()])
        except Exception as e:
            return jsonify({"error": str(e)}), 500

    @app.route("/api/files/convert", methods=["POST"])
    def api_files_convert():
        body = request.get_json(silent=True) or {}
        filename = str(body.get("filename", "")).strip()
        if not filename:
            return jsonify({"ok": False, "msg": "missing filename"}), 400
        src = Path(WECHAT_FILE_DIR) / filename
        if not src.exists():
            return jsonify({"ok": False, "msg": "file not found"}), 404
        def _do_convert():
            try:
                _convert_and_save(src)
            except Exception as e:
                LOG.error(f"[WebUI] 转换失败: {e}")
        Thread(target=_do_convert, daemon=True).start()
        return jsonify({"ok": True, "msg": f"已开始转换 {filename}"})

    @app.route("/api/files/converted/<path:name>", methods=["DELETE"])
    def api_delete_converted(name):
        p = FILE_SAVE_DIR / name
        if not p.exists():
            return jsonify({"error": "not found"}), 404
        p.unlink()
        return jsonify({"ok": True})

    # ── 日程编辑 ────────────────────────────────────────────────────────
    @app.route("/api/schedules", methods=["POST"])
    def api_create_schedule():
        body = request.get_json(silent=True) or {}
        name    = str(body.get("name", "")).strip()
        content = str(body.get("content", "")).strip()
        year    = int(body.get("year",  0))
        month   = int(body.get("month", 0))
        day     = int(body.get("day",   0))
        hour    = int(body.get("hour",  0))
        minute  = int(body.get("minute",0))
        weekday = body.get("weekday")  # None or 0-6
        if weekday is not None:
            weekday = int(weekday)
        if not name:
            return jsonify({"ok": False, "msg": "name is required"}), 400
        try:
            conn = sqlite3.connect(str(SCHEDULES_DB))
            conn.execute(
                "INSERT INTO schedules (name,year,month,day,hour,minute,weekday,content) VALUES (?,?,?,?,?,?,?,?)",
                (name, year, month, day, hour, minute, weekday, content)
            )
            conn.commit()
            new_id = conn.execute("SELECT last_insert_rowid()").fetchone()[0]
            conn.close()
            return jsonify({"ok": True, "id": new_id})
        except Exception as e:
            return jsonify({"error": str(e)}), 500

    @app.route("/api/schedules/<int:sid>", methods=["PUT"])
    def api_update_schedule(sid):
        body = request.get_json(silent=True) or {}
        name    = body.get("name")
        content = body.get("content")
        try:
            conn = sqlite3.connect(str(SCHEDULES_DB))
            if name is not None:
                conn.execute("UPDATE schedules SET name=? WHERE id=?", (name, sid))
            if content is not None:
                conn.execute("UPDATE schedules SET content=? WHERE id=?", (content, sid))
            conn.commit()
            conn.close()
            return jsonify({"ok": True})
        except Exception as e:
            return jsonify({"error": str(e)}), 500

    # ── 笔记编辑 ────────────────────────────────────────────────────────
    @app.route("/api/notes/<name>", methods=["PUT"])
    def api_update_note(name):
        body = request.get_json(silent=True) or {}
        content = body.get("content", "")
        matches = [p for p in NOTES_DIR.glob("*.txt") if p.stem.lower() == name.lower()]
        if not matches:
            return jsonify({"error": "not found"}), 404
        matches[0].write_text(content, encoding="utf-8")
        return jsonify({"ok": True})

    # ── 模型接口 ─────────────────────────────────────────────────────────
    @app.route("/api/model/provider", methods=["GET", "POST"])
    def api_model_provider():
        global AI_PROVIDER, AI_MODEL, client, _api_key, _base_url
        if request.method == "POST":
            body = request.get_json(silent=True) or {}
            p = str(body.get("provider", "")).lower().strip()
            if p not in _PROVIDER_CONFIGS:
                return jsonify({"ok": False, "msg": f"未知供应商: {p}"}), 400
            AI_PROVIDER = p
            _api_key, _base_url, AI_MODEL = _PROVIDER_CONFIGS[p]
            client = OpenAI(api_key=_api_key, base_url=_base_url)
            LOG.info(f"[WebUI] AI 供应商切换为: {AI_PROVIDER} ({AI_MODEL})")
            return jsonify({"ok": True, "provider": AI_PROVIDER, "model": AI_MODEL})
        return jsonify({"provider": AI_PROVIDER, "model": AI_MODEL, "available": list(_PROVIDER_CONFIGS.keys())})

    @app.route("/api/model/keys", methods=["GET", "POST"])
    def api_model_keys():
        global _PROVIDER_CONFIGS, _api_key, _base_url, AI_MODEL, client
        if request.method == "POST":
            body = request.get_json(silent=True) or {}
            provider = str(body.get("provider", "")).lower().strip()
            new_key  = str(body.get("key", "")).strip()
            if provider not in _PROVIDER_CONFIGS:
                return jsonify({"ok": False, "msg": "未知供应商"}), 400
            old_url, old_model = _PROVIDER_CONFIGS[provider][1], _PROVIDER_CONFIGS[provider][2]
            _PROVIDER_CONFIGS[provider] = (new_key, old_url, old_model)
            if provider == AI_PROVIDER:
                _api_key = new_key
                client = OpenAI(api_key=_api_key, base_url=_base_url)
            _try_write_env(f"{provider.upper()}_API_KEY", new_key)
            return jsonify({"ok": True})
        keys_info = {}
        for p, (k, url, model) in _PROVIDER_CONFIGS.items():
            masked = ("*" * (len(k) - 4) + k[-4:]) if k and len(k) > 4 else ("***" if k else "")
            keys_info[p] = {"masked": masked, "set": bool(k), "model": model, "base_url": url}
        return jsonify(keys_info)

    # ── Memory 管理 ──────────────────────────────────────────────────────
    @app.route("/api/memory")
    def api_memory_list():
        disabled = memory_disabled_load()
        files = sorted(MEMORY_DIR.glob("*.txt"), key=lambda p: p.stat().st_mtime, reverse=True)
        result = []
        for f in files:
            result.append({
                "name": f.stem,
                "filename": f.name,
                "size": f.stat().st_size,
                "mtime": int(f.stat().st_mtime),
                "enabled": f.name not in disabled,
            })
        return jsonify(result)

    @app.route("/api/memory/<name>")
    def api_memory_read(name):
        safe = re.sub(r'[\\/:*?"<>|]', '_', name)
        if not safe.endswith(".txt"):
            safe += ".txt"
        p = MEMORY_DIR / safe
        if not p.exists():
            return jsonify({"error": "not found"}), 404
        return jsonify({"name": p.stem, "filename": p.name, "content": p.read_text(encoding="utf-8", errors="ignore")})

    @app.route("/api/memory/<name>", methods=["PUT"])
    def api_memory_update(name):
        body = request.get_json(silent=True) or {}
        content = body.get("content", "")
        safe = re.sub(r'[\\/:*?"<>|]', '_', name)
        if not safe.endswith(".txt"):
            safe += ".txt"
        p = MEMORY_DIR / safe
        if not p.exists():
            return jsonify({"error": "not found"}), 404
        p.write_text(content, encoding="utf-8")
        return jsonify({"ok": True})

    @app.route("/api/memory/<name>", methods=["DELETE"])
    def api_memory_delete(name):
        ok = memory_delete(name)
        return jsonify({"ok": ok}) if ok else (jsonify({"error": "not found"}), 404)

    @app.route("/api/memory/<name>/toggle", methods=["POST"])
    def api_memory_toggle(name):
        safe = re.sub(r'[\\/:*?"<>|]', '_', name)
        if not safe.endswith(".txt"):
            safe += ".txt"
        is_disabled = memory_disabled_toggle(safe)
        return jsonify({"ok": True, "enabled": not is_disabled})

    # ── 网络配置（Sub Recorder） ─────────────────────────────────────────
    @app.route("/api/network/config")
    def api_network_config():
        return jsonify({"port": PUSH_PORT, "token_set": bool(PUSH_TOKEN)})

    @app.route("/api/network/token", methods=["POST"])
    def api_network_token():
        global PUSH_TOKEN
        body = request.get_json(silent=True) or {}
        new_token = str(body.get("token", "")).strip()
        PUSH_TOKEN = new_token
        _try_write_env("PUSH_TOKEN", new_token)
        return jsonify({"ok": True})

    @app.route("/api/push_test", methods=["POST"])
    def api_push_test():
        body = request.get_json(silent=True) or {}
        to = str(body.get("to", "")).strip()
        if not to:
            return jsonify({"ok": False, "msg": "missing to"}), 400
        q.put((to, "[WeChatBot] 推送测试消息 ✓"))
        return jsonify({"ok": True})

    LOG.info(f"[WebUI] 管理界面已启动: http://127.0.0.1:{PUSH_PORT}")
    app.run(host="0.0.0.0", port=PUSH_PORT, debug=False, use_reloader=False)


if __name__ == "__main__":
    main()
